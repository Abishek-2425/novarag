# ingest_pptx.py
import os
import logging
from typing import Dict, List, Optional
from pptx import Presentation

logging.basicConfig(level=logging.INFO, format="%(asctime)s [%(levelname)s] %(message)s")


class PPTXIngestor:
    """
    Handles ingestion of PPTX files and extraction of slide text for embeddings.
    Mirrors PDF/DOCX ingestors for seamless integration into multimodal RAG.
    """

    def __init__(self, folder_path: Optional[str] = None, chunk_size: int = 1000):
        """
        :param folder_path: Folder containing PPTX files
        :param chunk_size: Number of characters per text chunk
        """
        self.folder_path = folder_path
        self.chunk_size = chunk_size
        if self.folder_path and not os.path.exists(self.folder_path):
            raise FileNotFoundError(f"PPTX folder not found: {self.folder_path}")

    def extract_text_from_pptx(self, file_path: str) -> str:
        """
        Extracts all text from a PowerPoint (.pptx) file slide by slide.
        """
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"PPTX file not found: {file_path}")

        logging.info(f"Ingesting PPTX: {file_path}")
        text = ""
        try:
            prs = Presentation(file_path)
            for slide_num, slide in enumerate(prs.slides, start=1):
                slide_text = []
                for shape in slide.shapes:
                    # standard text shapes
                    if hasattr(shape, "text") and shape.text:
                        slide_text.append(shape.text)
                    # handle table cells separately
                    elif shape.shape_type == 19:  # MSO_SHAPE_TYPE.TABLE = 19
                        for row in shape.table.rows:
                            for cell in row.cells:
                                if cell.text.strip():
                                    slide_text.append(cell.text.strip())
                if slide_text:
                    text += f"\n--- Slide {slide_num} ---\n" + "\n".join(slide_text)
        except Exception as e:
            logging.error(f"Failed to extract text from {file_path}: {e}")
            return ""
        return text.strip()

    def chunk_text(self, text: str) -> List[str]:
        """
        Splits text into smaller chunks for embedding.
        """
        return [text[i:i + self.chunk_size] for i in range(0, len(text), self.chunk_size)]

    def ingest_file(self, file_path: str) -> List[Dict[str, str]]:
        """
        Ingest a single PPTX file and return list of dicts for each chunk:
        [{ 'text': chunk, 'source': filename }]
        """
        text = self.extract_text_from_pptx(file_path)
        chunks = self.chunk_text(text)
        return [{"text": chunk, "source": os.path.basename(file_path)} for chunk in chunks]

    def ingest_folder(self) -> List[Dict[str, str]]:
        """
        Extracts text from all PPTX files in the folder and returns list of dicts.
        Each dict contains a text chunk and the source filename.
        """
        if not self.folder_path:
            raise ValueError("Folder path not set for folder ingestion.")

        all_texts = []
        for filename in os.listdir(self.folder_path):
            if filename.lower().endswith(".pptx"):
                full_path = os.path.join(self.folder_path, filename)
                try:
                    file_chunks = self.ingest_file(full_path)
                    all_texts.extend(file_chunks)
                except Exception as e:
                    logging.error(f"Failed to ingest {filename}: {e}")
        logging.info(f"Total PPTX chunks ingested: {len(all_texts)}")
        return all_texts


# Example usage
if __name__ == "__main__":
    folder = "sample_data/pptx"
    pptx_ingestor = PPTXIngestor(folder_path=folder)
    chunks = pptx_ingestor.ingest_folder()
    for i, item in enumerate(chunks):
        print(f"----- Chunk {i+1} from {item['source']} -----\n{item['text'][:500]}...\n")
